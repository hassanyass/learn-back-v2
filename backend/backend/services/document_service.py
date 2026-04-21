import io
import os
import subprocess
from pathlib import Path

from fastapi import HTTPException, UploadFile, status
from PyPDF2 import PdfReader
from pptx import Presentation


class DocumentService:
    ALLOWED_EXTENSIONS = {".pdf", ".pptx"}

    def validate_upload(self, upload_file: UploadFile) -> str:
        suffix = Path(upload_file.filename or "").suffix.lower()
        if suffix not in self.ALLOWED_EXTENSIONS:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Only PDF and PPTX files are supported.",
            )
        return suffix

    async def extract_raw_text(self, upload_file: UploadFile) -> str:
        extension = self.validate_upload(upload_file)
        content = await upload_file.read()
        if not content:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Uploaded file is empty.",
            )

        if extension == ".pdf":
            return self._extract_pdf_text(content)
        return self._extract_pptx_text(content)

    def _extract_pdf_text(self, content: bytes) -> str:
        reader = PdfReader(io.BytesIO(content))
        pages = [(page.extract_text() or "").strip() for page in reader.pages]
        text = "\n".join(filter(None, pages)).strip()
        if not text:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail="No extractable text found in PDF.",
            )
        return text

    def _extract_pptx_text(self, content: bytes) -> str:
        prs = Presentation(io.BytesIO(content))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text.strip())
        text = "\n".join(filter(None, lines)).strip()
        if not text:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail="No extractable text found in PPTX.",
            )
        return text

    def convert_pptx_to_pdf(self, input_pptx_path: str, output_dir: str) -> str:
        """
        Stub for server-side PPTX->PDF conversion.
        Requires LibreOffice in runtime image.
        """
        command = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            output_dir,
            input_pptx_path,
        ]
        subprocess.run(command, check=True)
        output_pdf = Path(output_dir) / f"{Path(input_pptx_path).stem}.pdf"
        return str(output_pdf)

    def upload_pdf_to_storage(self, pdf_path: str, storage_key: str) -> str:
        """
        Stub for Supabase Storage upload.
        Replace with official Supabase client integration.
        """
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        return f"https://example.supabase.co/storage/v1/object/public/slides/{storage_key}"
